"""
reports.pkg에 신규/업데이트 보고서를 증분 추가하는 스크립트
- 디스크 기반 ZIP + AES-CTR 스트리밍 암호화로 대용량 패키지 지원
- v1(Fernet) 패키지를 읽어 v2(AES-CTR) 패키지로 저장

사용법: python add_reports.py [현장코드1 현장코드2 ...]
  - 인자 없이 실행하면 Reports/ 전체를 기준으로 패키지 동기화
  - 현장코드를 지정하면 해당 코드만 업데이트
"""
import os
import io
import sys
import json
import struct
import hmac as hmac_mod
import hashlib
import zipfile
import base64
import shutil
import tempfile
from pathlib import Path
from cryptography.fernet import Fernet
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
import fitz  # PyMuPDF

from config import BASE_DIR, REPORTS_ORIGINAL_DIR, PACKAGE_FILE, MASTER_PASSWORD

BUILD_DIR = BASE_DIR / 'build_temp_incr'
CHUNK_SIZE = 64 * 1024 * 1024  # 64MB


def derive_fernet_key(password: str, salt: bytes) -> bytes:
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(), length=32,
        salt=salt, iterations=480000,
    )
    return base64.urlsafe_b64encode(kdf.derive(password.encode()))


def derive_raw_key(password: str, salt: bytes) -> bytes:
    """AES-CTR용 32바이트 원시 키"""
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(), length=32,
        salt=salt, iterations=480000,
    )
    return kdf.derive(password.encode())


def load_existing_package():
    """기존 패키지 로드 → (ZipFile, manifest) 반환"""
    if not PACKAGE_FILE.exists():
        return None, {'version': '2.0', 'reports': {}}

    with open(PACKAGE_FILE, 'rb') as f:
        magic = f.read(8)

    if magic == b'CJPHv1\x00\x00':
        return _load_v1()
    elif magic == b'CJPHv2\x00\x00':
        return _load_v2()
    else:
        raise ValueError("유효하지 않은 패키지 파일")


def _load_v1():
    """Fernet 방식 (기존)"""
    with open(PACKAGE_FILE, 'rb') as f:
        f.read(8)
        salt = f.read(16)
        encrypted = f.read()

    key = derive_fernet_key(MASTER_PASSWORD, salt)
    cipher = Fernet(key)
    decrypted = cipher.decrypt(encrypted)

    old_zip = zipfile.ZipFile(io.BytesIO(decrypted), 'r')
    manifest = json.loads(old_zip.read('manifest.json'))
    return old_zip, manifest


def _load_v2():
    """AES-CTR 스트리밍 방식"""
    with open(PACKAGE_FILE, 'rb') as f:
        f.read(8)
        salt = f.read(16)
        nonce = f.read(16)
        stored_mac = f.read(32)
        encrypted_data = f.read()

    raw_key = derive_raw_key(MASTER_PASSWORD, salt)
    computed_mac = hmac_mod.new(raw_key, encrypted_data, hashlib.sha256).digest()
    if not hmac_mod.compare_digest(stored_mac, computed_mac):
        raise ValueError("HMAC 검증 실패")

    cipher = Cipher(algorithms.AES(raw_key), modes.CTR(nonce))
    decryptor = cipher.decryptor()

    buf = io.BytesIO()
    offset = 0
    while offset < len(encrypted_data):
        chunk = encrypted_data[offset:offset + CHUNK_SIZE]
        buf.write(decryptor.update(chunk))
        offset += CHUNK_SIZE
    buf.write(decryptor.finalize())
    buf.seek(0)

    old_zip = zipfile.ZipFile(buf, 'r')
    manifest = json.loads(old_zip.read('manifest.json'))
    return old_zip, manifest


def convert_pdf_to_images(pdf_path: Path, out_dir: Path):
    """PDF → PNG (PyMuPDF)"""
    out_dir.mkdir(parents=True, exist_ok=True)
    images = []
    doc = fitz.open(str(pdf_path))
    for i, page in enumerate(doc, 1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = out_dir / f"page_{i:03d}.png"
        pix.save(str(img_path))
        images.append(img_path)
    doc.close()
    return images


def convert_pptx_to_images(pptx_path: Path, out_dir: Path):
    """PowerPoint COM 자동화로 PPT → PNG"""
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    powerpoint = None
    deck = None
    images = []

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(
            str(pptx_path.absolute()), WithWindow=False
        )
        out_dir.mkdir(parents=True, exist_ok=True)

        for i, slide in enumerate(deck.Slides, 1):
            img_path = out_dir / f"slide_{i:03d}.png"
            slide.Export(str(img_path.absolute()), "PNG", 1920, 1080)
            images.append(img_path)
    finally:
        if deck:
            deck.Close()
        if powerpoint:
            powerpoint.Quit()
        pythoncom.CoUninitialize()

    return images


def extract_text(filepath: Path):
    """검색 인덱싱용 페이지별 텍스트 추출"""
    suffix = filepath.suffix.lower()
    pages = []
    try:
        if suffix in ['.pptx', '.ppt']:
            from pptx import Presentation
            prs = Presentation(filepath)
            for idx, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                texts.append(cell.text)
                content = ' '.join(t.strip() for t in texts if t and t.strip())
                if content:
                    pages.append((idx, content))
        elif suffix == '.pdf':
            doc = fitz.open(str(filepath))
            for idx, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    pages.append((idx, text))
            doc.close()
    except Exception as e:
        print(f"   [경고] 텍스트 추출 실패: {e}")
    return pages


def encrypt_file_streaming(input_path: Path, output_path: Path, password: str):
    """
    AES-CTR 스트리밍 암호화 → CJPHv2 형식
    메모리 사용량: CHUNK_SIZE (64MB) 수준만 사용
    """
    salt = os.urandom(16)
    nonce = os.urandom(16)
    raw_key = derive_raw_key(password, salt)

    cipher = Cipher(algorithms.AES(raw_key), modes.CTR(nonce))
    encryptor = cipher.encryptor()

    # 1단계: 암호화된 데이터를 임시 파일에 기록
    temp_enc = output_path.with_suffix('.tmp_enc')
    mac = hmac_mod.new(raw_key, digestmod=hashlib.sha256)

    with open(input_path, 'rb') as fin, open(temp_enc, 'wb') as fout:
        while True:
            chunk = fin.read(CHUNK_SIZE)
            if not chunk:
                break
            encrypted_chunk = encryptor.update(chunk)
            mac.update(encrypted_chunk)
            fout.write(encrypted_chunk)
        final = encryptor.finalize()
        if final:
            mac.update(final)
            fout.write(final)

    hmac_digest = mac.digest()

    # 2단계: 최종 .pkg 파일 조립
    with open(output_path, 'wb') as fout:
        fout.write(b'CJPHv2\x00\x00')  # magic (8)
        fout.write(salt)                 # salt (16)
        fout.write(nonce)                # nonce (16)
        fout.write(hmac_digest)          # HMAC (32)

        # 암호화된 데이터 복사
        with open(temp_enc, 'rb') as fin:
            while True:
                chunk = fin.read(CHUNK_SIZE)
                if not chunk:
                    break
                fout.write(chunk)

    # 임시 파일 삭제
    temp_enc.unlink(missing_ok=True)


def add_reports(target_codes=None):
    """
    target_codes: 업데이트할 현장코드 리스트. None이면 패키지에 없는 것만 추가.
    """
    print("=" * 65)
    print("  CJ ProjectHub - 보고서 증분 업데이트 (v2 스트리밍)")
    print("=" * 65)

    if not REPORTS_ORIGINAL_DIR.exists():
        print(f"\n[오류] Reports 폴더 없음: {REPORTS_ORIGINAL_DIR}")
        return

    # 1. 기존 패키지 로드
    print(f"\n[1/6] 기존 패키지 로드...")
    try:
        old_zip, manifest = load_existing_package()
        print(f"   기존 보고서: {len(manifest['reports'])}건")
    except Exception as e:
        print(f"[오류] 패키지 로드 실패: {e}")
        print("   새 패키지를 생성합니다.")
        old_zip = None
        manifest = {'version': '2.0', 'reports': {}}

    # 2. 대상 파일 결정
    all_files = {}
    for f in REPORTS_ORIGINAL_DIR.glob('*'):
        if f.suffix.lower() in ['.pptx', '.ppt', '.pdf', '.xls', '.xlsx']:
            code = f.stem.split('_')[0]
            if code.isdigit():
                all_files[code] = f

    if target_codes:
        process_codes = [c for c in target_codes if c in all_files]
        missing = [c for c in target_codes if c not in all_files]
        if missing:
            print(f"\n[경고] Reports/에 없는 코드: {missing}")
    else:
        process_codes = [c for c in all_files if c not in manifest['reports']]

    if not process_codes:
        print("\n[완료] 업데이트할 보고서가 없습니다.")
        return

    print(f"\n[2/6] 대상: {len(process_codes)}건")

    # 3. 빌드 디렉토리 준비
    print(f"[3/6] 빌드 디렉토리 준비...")
    if BUILD_DIR.exists():
        shutil.rmtree(BUILD_DIR)
    BUILD_DIR.mkdir(parents=True)

    # 기존 패키지에서 업데이트 대상이 아닌 파일만 추출
    if old_zip:
        skip_prefixes = set()
        for code in process_codes:
            skip_prefixes.add(f"{code}/")

        for name in old_zip.namelist():
            skip = False
            for prefix in skip_prefixes:
                if name.startswith(prefix):
                    skip = True
                    break
            if skip:
                continue  # 업데이트 대상은 건너뜀
            old_zip.extract(name, BUILD_DIR)
        old_zip.close()

    # 4. 대상 파일 처리
    print(f"[4/6] 보고서 변환 중...")
    success = 0
    fail = 0
    for idx, code in enumerate(process_codes, 1):
        filepath = all_files[code]
        site_dir = BUILD_DIR / code

        if site_dir.exists():
            shutil.rmtree(site_dir)
        site_dir.mkdir(parents=True)

        print(f"  [{idx}/{len(process_codes)}] {filepath.name}")

        try:
            suffix = filepath.suffix.lower()
            if suffix in ['.pptx', '.ppt']:
                images = convert_pptx_to_images(filepath, site_dir)
            elif suffix == '.pdf':
                images = convert_pdf_to_images(filepath, site_dir)
            else:
                print(f"   [SKIP] 지원하지 않는 형식: {suffix}")
                continue

            text_pages = extract_text(filepath)

            manifest['reports'][code] = {
                'original_name': filepath.name,
                'page_count': len(images),
                'text_index': [{'page': p, 'content': t} for p, t in text_pages]
            }
            print(f"   ✓ {len(images)}p 변환, {len(text_pages)}p 텍스트")
            success += 1

        except Exception as e:
            print(f"   ✗ 실패: {e}")
            fail += 1

    # 5. manifest 저장 & 디스크 기반 ZIP 생성
    print(f"\n[5/6] ZIP 압축 중 (디스크 기반)...")
    manifest['version'] = '2.0'
    with open(BUILD_DIR / 'manifest.json', 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    zip_path = BASE_DIR / 'build_temp.zip'
    with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in BUILD_DIR.rglob('*'):
            if f.is_file():
                zf.write(f, f.relative_to(BUILD_DIR))

    zip_size = zip_path.stat().st_size
    print(f"   ZIP 크기: {zip_size / 1024 / 1024:.1f} MB")

    # 빌드 디렉토리 정리 (ZIP 생성 후 즉시)
    shutil.rmtree(BUILD_DIR)

    # 6. 스트리밍 암호화
    print(f"[6/6] 스트리밍 암호화 중...")
    encrypt_file_streaming(zip_path, PACKAGE_FILE, MASTER_PASSWORD)

    # 임시 ZIP 삭제
    zip_path.unlink(missing_ok=True)

    pkg_size = PACKAGE_FILE.stat().st_size
    print("\n" + "=" * 65)
    print(f"  ✓ 패키지 업데이트 완료: {PACKAGE_FILE.name}")
    print(f"  ✓ 총 보고서: {len(manifest['reports'])}건")
    print(f"  ✓ 이번 처리: 성공 {success}건 / 실패 {fail}건")
    print(f"  ✓ 패키지 크기: {pkg_size / 1024 / 1024:.1f} MB")
    print(f"  ✓ 형식: CJPHv2 (AES-CTR 스트리밍)")
    print("=" * 65)


if __name__ == '__main__':
    codes = sys.argv[1:] if len(sys.argv) > 1 else None
    if codes:
        print(f"지정 코드 업데이트: {codes}")
    add_reports(codes)
