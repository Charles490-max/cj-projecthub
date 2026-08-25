"""
원본 PPT/PDF를 이미지로 변환하고 암호화하여 reports.pkg 생성
※ Charles님이 신규 보고서 추가 시 1회 실행
"""
import os
import io
import json
import zipfile
import shutil
import base64
from pathlib import Path
from cryptography.fernet import Fernet
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.primitives import hashes
import fitz  # PyMuPDF

from config import (
    BASE_DIR, REPORTS_ORIGINAL_DIR, PACKAGE_FILE, MASTER_PASSWORD
)

BUILD_DIR = BASE_DIR / 'build_temp'


def derive_key(password: str, salt: bytes) -> bytes:
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(), length=32,
        salt=salt, iterations=480000,
    )
    return base64.urlsafe_b64encode(kdf.derive(password.encode()))


def convert_pptx_to_images(pptx_path: Path, out_dir: Path):
    """PowerPoint COM 자동화로 PPT → PNG"""
    import win32com.client
    import pythoncom
    
    pythoncom.CoInitialize()
    powerpoint = None
    deck = None
    images = []
    
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        
        deck = powerpoint.Presentations.Open(
            str(pptx_path.absolute()), WithWindow=False
        )
        out_dir.mkdir(parents=True, exist_ok=True)
        
        for i, slide in enumerate(deck.Slides, 1):
            img_path = out_dir / f"slide_{i:03d}.png"
            slide.Export(str(img_path.absolute()), "PNG", 1920, 1080)
            images.append(img_path)
    
    finally:
        if deck:
            deck.Close()
        if powerpoint:
            powerpoint.Quit()
        pythoncom.CoUninitialize()
    
    return images


def convert_pdf_to_images(pdf_path: Path, out_dir: Path):
    """PDF → PNG (PyMuPDF)"""
    out_dir.mkdir(parents=True, exist_ok=True)
    images = []
    doc = fitz.open(str(pdf_path))
    for i, page in enumerate(doc, 1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = out_dir / f"page_{i:03d}.png"
        pix.save(str(img_path))
        images.append(img_path)
    doc.close()
    return images


def extract_text(filepath: Path):
    """검색 인덱싱용 페이지별 텍스트 추출"""
    suffix = filepath.suffix.lower()
    pages = []
    
    try:
        if suffix in ['.pptx', '.ppt']:
            from pptx import Presentation
            prs = Presentation(filepath)
            for idx, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                texts.append(cell.text)
                content = ' '.join(t.strip() for t in texts if t and t.strip())
                if content:
                    pages.append((idx, content))
        
        elif suffix == '.pdf':
            doc = fitz.open(str(filepath))
            for idx, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    pages.append((idx, text))
            doc.close()
    except Exception as e:
        print(f"   [경고] 텍스트 추출 실패: {e}")
    
    return pages


def build_package():
    print("=" * 65)
    print(f"  CJ ProjectHub - 보안 패키지 빌드")
    print("=" * 65)
    
    if not REPORTS_ORIGINAL_DIR.exists():
        print(f"\n[오류] 원본 폴더 없음: {REPORTS_ORIGINAL_DIR}")
        return
    
    # 빌드 임시 폴더 초기화
    if BUILD_DIR.exists():
        shutil.rmtree(BUILD_DIR)
    BUILD_DIR.mkdir(parents=True)
    
    manifest = {'version': '1.0', 'reports': {}}
    success_count = 0
    fail_count = 0
    
    # 지원 파일만 수집
    target_files = []
    for f in REPORTS_ORIGINAL_DIR.glob('*'):
        if f.suffix.lower() in ['.pptx', '.ppt', '.pdf']:
            target_files.append(f)
    
    print(f"\n[스캔] 처리 대상 파일: {len(target_files)}건\n")
    
    for idx, orig_file in enumerate(target_files, 1):
        # 현장코드 추출
        site_code = orig_file.stem.split('_')[0]
        if not site_code.isdigit():
            print(f"[{idx}/{len(target_files)}] [SKIP] 코드 없음: {orig_file.name}")
            continue
        
        # 동일 site_code 중복 처리 (PPT, PDF 둘 다 있는 경우 PPT 우선)
        if site_code in manifest['reports']:
            print(f"[{idx}/{len(target_files)}] [SKIP] 중복: {orig_file.name}")
            continue
        
        print(f"[{idx}/{len(target_files)}] {orig_file.name}")
        site_dir = BUILD_DIR / site_code
        
        try:
            if orig_file.suffix.lower() in ['.pptx', '.ppt']:
                images = convert_pptx_to_images(orig_file, site_dir)
            else:
                images = convert_pdf_to_images(orig_file, site_dir)
            
            text_pages = extract_text(orig_file)
            
            manifest['reports'][site_code] = {
                'original_name': orig_file.name,
                'page_count': len(images),
                'text_index': [{'page': p, 'content': t} for p, t in text_pages]
            }
            print(f"   ✓ {len(images)}페이지 변환 완료")
            success_count += 1
        
        except Exception as e:
            print(f"   ✗ 실패: {e}")
            fail_count += 1
    
    # manifest 저장
    with open(BUILD_DIR / 'manifest.json', 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    
    # ZIP 압축 (디스크 기반)
    print(f"\n[압축 중...]")
    zip_path = BASE_DIR / 'build_temp.zip'
    with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in BUILD_DIR.rglob('*'):
            if f.is_file():
                zf.write(f, f.relative_to(BUILD_DIR))
    zip_size = zip_path.stat().st_size
    print(f"   ZIP 크기: {zip_size/1024/1024:.1f} MB")

    # 임시 폴더 정리
    shutil.rmtree(BUILD_DIR)

    # AES-CTR 스트리밍 암호화 (v2)
    print(f"[스트리밍 암호화 중...]")
    from add_reports import encrypt_file_streaming
    encrypt_file_streaming(zip_path, PACKAGE_FILE, MASTER_PASSWORD)

    # 임시 ZIP 삭제
    zip_path.unlink(missing_ok=True)

    print("\n" + "=" * 65)
    print(f"  ✓ 완료: {PACKAGE_FILE.name}")
    print(f"  ✓ 패키지 크기: {PACKAGE_FILE.stat().st_size/1024/1024:.1f} MB")
    print(f"  ✓ 성공 {success_count}건 / 실패 {fail_count}건")
    print(f"  ✓ 형식: CJPHv2 (AES-CTR 스트리밍)")
    print("=" * 65)


if __name__ == '__main__':
    build_package()
